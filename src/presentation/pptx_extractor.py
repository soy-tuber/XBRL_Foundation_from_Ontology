"""
PPTX → スライド単位テキスト抽出。python-pptx を使用。
"""

from __future__ import annotations

from typing import List

from src.presentation.pdf_extractor import Slide


def extract_pptx(path: str) -> List[Slide]:
    try:
        from pptx import Presentation
    except ImportError as e:
        raise RuntimeError("python-pptx が必要です: pip install python-pptx") from e

    prs = Presentation(path)
    slides: List[Slide] = []
    for i, slide in enumerate(prs.slides, start=1):
        title = None
        body_parts: List[str] = []
        has_table = False
        has_chart = False

        for shape in slide.shapes:
            if shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in p.runs).strip()
                    if line:
                        if title is None and shape == slide.shapes.title:
                            title = line
                        body_parts.append(line)
            if getattr(shape, "has_table", False):
                has_table = True
                try:
                    tbl = shape.table
                    for row in tbl.rows:
                        body_parts.append("\t".join(cell.text.strip() for cell in row.cells))
                except Exception:
                    pass
            if getattr(shape, "has_chart", False):
                has_chart = True

        body = "\n".join(body_parts).strip()
        if not body and not (has_table or has_chart):
            continue
        slides.append(Slide(
            slide_no=i,
            title=title,
            content_text=body,
            has_table=has_table,
            has_chart=has_chart,
        ))
    return slides
